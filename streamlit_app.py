import os
import io
import csv
from dotenv import load_dotenv
load_dotenv()

from pathlib import Path
from dataclasses import dataclass
from typing import List, Tuple

import numpy as np
import streamlit as st
from openai import OpenAI
from pypdf import PdfReader
from pptx import Presentation
from openpyxl import load_workbook


# -----------------------------
# Text extraction (from bytes)
# -----------------------------

def extract_pptx_bytes(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_parts: List[str] = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                txt = (shape.text or "").strip()
                if txt:
                    slide_parts.append(txt)
        # Notes
        try:
            notes = slide.notes_slide
            if notes and notes.notes_text_frame:
                ntext = (notes.notes_text_frame.text or "").strip()
                if ntext:
                    slide_parts.append(f"[Notes] {ntext}")
        except Exception:
            pass
        parts.append("\n".join(slide_parts))
    return "\n\n".join(parts).strip()

def extract_pdf_bytes(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    parts: List[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = (page.extract_text() or "").strip()
        except Exception:
            txt = ""
        if txt:
            parts.append(f"[Page {i}]\n{txt}")
    return "\n\n".join(parts).strip()

def extract_txt_bytes(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore").strip()

def extract_csv_bytes(data: bytes, max_rows: int = 5000) -> str:
    s = data.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(s))
    parts: List[str] = []
    for i, row in enumerate(reader, start=1):
        if i > max_rows:
            parts.append(f"[Truncated after {max_rows} rows]")
            break
        line = " | ".join((c or "").strip() for c in row)
        if line.strip():
            parts.append(line)
    return "\n".join(parts).strip()

def extract_xlsx_bytes(data: bytes, max_cells: int = 20000) -> str:
    # openpyxl needs a file-like object
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts: List[str] = []
    cell_count = 0

    for ws in wb.worksheets:
        parts.append(f"[Sheet] {ws.title}")
        for row in ws.iter_rows(values_only=True):
            if cell_count >= max_cells:
                parts.append(f"[Truncated after {max_cells} cells]")
                wb.close()
                return "\n".join(parts).strip()

            rendered: List[str] = []
            for v in row:
                cell_count += 1
                rendered.append("" if v is None else str(v).strip())

            if any(x for x in rendered):
                parts.append(" | ".join(rendered))
        parts.append("")
    wb.close()
    return "\n".join(parts).strip()

def load_uploaded_files(files) -> str:
    supported = {".pptx", ".pdf", ".txt", ".md", ".csv", ".xlsx"}
    outputs: List[str] = []

    for uf in files:
        name = uf.name
        ext = Path(name).suffix.lower()
        if ext not in supported:
            raise ValueError(f"Unsupported file type: {ext} ({name})")

        data = uf.getvalue()
        header = f"[Source {name}]"

        if ext == ".pptx":
            body = extract_pptx_bytes(data)
        elif ext == ".pdf":
            body = extract_pdf_bytes(data)
        elif ext in {".txt", ".md"}:
            body = extract_txt_bytes(data)
        elif ext == ".csv":
            body = extract_csv_bytes(data)
        elif ext == ".xlsx":
            body = extract_xlsx_bytes(data)
        else:
            body = ""

        if body:
            outputs.append(f"{header}\n{body}")

    return "\n\n=====\n\n".join(outputs).strip()


# -----------------------------
# Chunking + vector search
# -----------------------------

def chunk_text(text: str, max_chars: int = 1200, overlap: int = 150) -> List[str]:
    if not text:
        return []
    chunks: List[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + max_chars, n)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end == n:
            break
        start = max(0, end - overlap)
    return chunks

@dataclass
class DocChunk:
    text: str
    embedding: np.ndarray

class InMemoryVectorStore:
    def __init__(self, chunks: List[DocChunk]):
        self.chunks = chunks
        if chunks:
            m = np.vstack([c.embedding for c in chunks]).astype(np.float32)
            norms = np.linalg.norm(m, axis=1, keepdims=True)
            self.matrix = m / np.clip(norms, 1e-12, None)
        else:
            self.matrix = np.zeros((0, 1), dtype=np.float32)

    def top_k(self, query_embedding: np.ndarray, k: int = 5) -> List[Tuple[float, str]]:
        if self.matrix.shape[0] == 0:
            return []
        q = query_embedding.astype(np.float32)
        q = q / max(np.linalg.norm(q), 1e-12)
        sims = self.matrix @ q
        idx = np.argsort(-sims)[:k]
        return [(float(sims[i]), self.chunks[int(i)].text) for i in idx]


# -----------------------------
# OpenAI
# -----------------------------

def get_client() -> OpenAI:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY is not set.")
    return OpenAI(api_key=api_key)

def embed_texts(client: OpenAI, texts: List[str]) -> List[np.ndarray]:
    resp = client.embeddings.create(
        model="text-embedding-3-small",
        input=texts,
    )
    return [np.array(d.embedding, dtype=np.float32) for d in resp.data]

def answer_with_context(client: OpenAI, question: str, context_chunks: List[str]) -> str:
    context = "\n\n---\n\n".join(context_chunks)
    system = (
        "You are HOPE, a helpful assistant. Answer using ONLY the provided context. "
        "If the answer isn't in the context, say you don't know and ask what document/page/slide to use. "
        "Cite slide/page markers like [Slide 3] or [Page 2] when present."
    )
    user = (
        "CONTEXT:\n"
        f"{context}\n\n"
        "QUESTION:\n"
        f"{question}\n"
    )
    resp = client.responses.create(
        model="gpt-4.1-mini",
        input=[
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        temperature=0.2,
    )
    return resp.output_text.strip()


# -----------------------------
# Streamlit UI
# -----------------------------

# Show logo + title side by side
ICON_PATH = Path(__file__).parent / "icon.png"
col1, col2 = st.columns([1,6])
with col1:
    if ICON_PATH.exists():
        st.image(str(ICON_PATH), width=70)
with col2:
    st.title("HOPE Chatbot")
st.caption("Upload PPTX/PDF/Excel/CSV/TXT and ask questions grounded in those documents.")

top_k = st.sidebar.slider("Context chunks (top_k)", 3, 12, 5)

uploaded = st.sidebar.file_uploader(
    "Upload documents",
    type=["pptx", "pdf", "txt", "md", "csv", "xlsx"],
    accept_multiple_files=True
)

index_btn = st.sidebar.button("Index documents")

if "store" not in st.session_state:
    st.session_state.store = None
if "messages" not in st.session_state:
    st.session_state.messages = []

if index_btn:
    if not uploaded:
        st.sidebar.error("Upload at least one file.")
    else:
        client = get_client()
        with st.spinner("Extracting text..."):
            corpus = load_uploaded_files(uploaded)
        chunks = chunk_text(corpus)

        with st.spinner(f"Embedding {len(chunks)} chunks..."):
            embs = embed_texts(client, chunks)

        doc_chunks = [DocChunk(t, e) for t, e in zip(chunks, embs)]
        st.session_state.store = InMemoryVectorStore(doc_chunks)
        st.sidebar.success(f"Indexed {len(chunks)} chunks from {len(uploaded)} file(s).")

for m in st.session_state.messages:
    with st.chat_message(m["role"]):
        st.markdown(m["content"])

prompt = st.chat_input("Ask a question about the uploaded documents...")

if prompt:
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    if st.session_state.store is None:
        with st.chat_message("assistant"):
            st.markdown("Upload and index documents first (left sidebar).")
    else:
        client = get_client()
        q_emb = embed_texts(client, [prompt])[0]
        hits = st.session_state.store.top_k(q_emb, k=top_k)
        context_chunks = [t for _, t in hits]

        answer = answer_with_context(client, prompt, context_chunks)
        st.session_state.messages.append({"role": "assistant", "content": answer})

        with st.chat_message("assistant"):
            st.markdown(answer)
            with st.expander("Retrieved context (debug)"):
                for score, txt in hits:
                    st.markdown(f"**score:** {score:.3f}\n\n{txt}")